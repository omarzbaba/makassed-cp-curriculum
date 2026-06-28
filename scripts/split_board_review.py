#!/usr/bin/env python3
"""Split board-review Q&A slides: question on one slide, answer revealed on the next.

For each board-review slide (question + options + an ANSWER panel + explanation all
on one slide), we duplicate the slide and then strip the answer panel from the FIRST
copy. Result: a clean question slide, followed by the same slide with the answer shown.

The decks contain no pictures on these slides (verified), so deep-copying shape XML
is lossless. Detection is geometric & relative to the ANSWER shape, so it is robust
to per-deck position shifts.
"""
import copy
import glob
import os
import re
import sys
from collections import Counter
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor

FOOTER_PREFIXES = ("Makassed General Hospital", "CP Lecture Series")
OPTION_RE = re.compile(r"^[A-E][.)]\s")


def normalize_options(slide):
    """On a question slide, reset the (pre-highlighted) correct option back to the
    same neutral styling as the other options, so the answer isn't given away."""
    opt, best = None, 0
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        n = sum(1 for p in sh.text_frame.paragraphs if OPTION_RE.match(p.text.strip()))
        if n > best:
            best, opt = n, sh
    if not opt or best < 2:
        return False
    cols = []
    for p in opt.text_frame.paragraphs:
        for r in p.runs:
            try:
                cols.append(str(r.font.color.rgb))
            except Exception:
                pass
    if not cols:
        return False
    normal = Counter(cols).most_common(1)[0][0]   # the body color shared by most options
    changed = False
    for p in opt.text_frame.paragraphs:
        for r in p.runs:
            try:
                cur = str(r.font.color.rgb)
            except Exception:
                cur = None
            if r.font.bold or (cur is not None and cur != normal):
                r.font.bold = False
                r.font.color.rgb = RGBColor.from_string(normal)
                changed = True
    return changed


def is_board_q(slide):
    txts = [sh.text_frame.text.strip() for sh in slide.shapes if sh.has_text_frame]
    joined = " ".join(txts).lower()
    return ("board" in joined) and any(t.upper().startswith("ANSWER") for t in txts)


def answer_panel_elements(slide):
    """Return the XML elements of the answer panel (bg + ANSWER + explanation)."""
    ans = next((sh for sh in slide.shapes
                if sh.has_text_frame and sh.text_frame.text.strip().upper().startswith("ANSWER")), None)
    if ans is None:
        return None
    ax, ay = ans.left, ans.top
    els = []
    for sh in slide.shapes:
        if sh.left is None or sh.top is None:
            continue
        t = sh.text_frame.text.strip() if sh.has_text_frame else ""
        if t.startswith(FOOTER_PREFIXES):
            continue
        if sh.left >= ax - Inches(0.6) and (ay - Inches(0.6)) <= sh.top <= (ay + Inches(3)):
            els.append(sh._element)
    # Safety: the ANSWER shape itself must be in the set
    if ans._element not in els:
        return None
    return els


def duplicate_slide_after(prs, index):
    src = prs.slides[index]
    dest = prs.slides.add_slide(src.slide_layout)
    # drop placeholders the layout injected
    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)
    # lossless copy of every shape (no pictures on these slides)
    for shp in src.shapes:
        dest.shapes._spTree.append(copy.deepcopy(shp._element))
    # move the just-added slide (currently last) to sit right after `index`
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    sldIdLst.remove(ids[-1])
    sldIdLst.insert(index + 1, ids[-1])
    return dest


def process_deck(path, out_path):
    prs = Presentation(path)
    targets = [i for i, sl in enumerate(prs.slides) if is_board_q(sl)]
    # process from last to first so earlier indices don't shift
    splits = 0
    for i in sorted(targets, reverse=True):
        els = answer_panel_elements(prs.slides[i])
        if not els:
            raise RuntimeError(f"{os.path.basename(path)} slide {i}: could not locate answer panel")
        duplicate_slide_after(prs, i)            # i+1 becomes the full answer slide
        for el in els:                            # strip answer from the question slide (i)
            el.getparent().remove(el)
        normalize_options(prs.slides[i])         # de-highlight the correct option on the question slide
        splits += 1
    prs.save(out_path)
    return splits, len(targets), len(Presentation(out_path).slides)


if __name__ == "__main__":
    files = sys.argv[1:] or sorted(glob.glob("decks/*.pptx"))
    for f in files:
        splits, targets, total = process_deck(f, f)  # in place
        print(f"{os.path.basename(f):55} split {splits:2} board Qs -> {total} slides")
